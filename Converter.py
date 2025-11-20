import pandas as pd
from docx import Document
from docx.shared import Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.enum.section import WD_SECTION
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.oxml.xmlchemy import OxmlElement
import argparse
import os
import sys
import config


# =========================================================
# VALIDATION
# =========================================================

def validate_excel(df):
    """Validates the structure and content of the Excel data."""
    print("\n🔍 Performing validation checks on input file...")

    # 1️⃣ Question No validation
    try:
        question_nos = df["Question No"].tolist()
    except KeyError:
        print("❌ 'Question No' column is missing.")
        sys.exit(1)

    valid_numbers = []
    invalid_rows = []
    for idx, val in enumerate(question_nos, start=2):
        try:
            num = int(str(val).strip())
            valid_numbers.append(num)
        except ValueError:
            invalid_rows.append(idx)

    if invalid_rows or len(valid_numbers) != len(set(valid_numbers)):
        print("❌ File is not compatible for conversion.")
        print("   All 'Question No' values must be numeric and appear exactly once.")
        if invalid_rows:
            print(f"   → Non-numeric values in Excel rows: {invalid_rows}")
        sys.exit(1)

    expected = set(range(1, len(valid_numbers) + 1))
    missing = sorted(list(expected - set(valid_numbers)))
    extra = sorted(list(set(valid_numbers) - expected))

    if missing or extra:
        print("❌ File is not compatible for conversion.")
        if missing:
            print(f"   → Missing Question Numbers: {missing}")
        if extra:
            print(f"   → Unexpected Question Numbers: {extra}")
        sys.exit(1)

    # 2️⃣ Answer column validation
    try:
        answers = df["answer"].tolist()
    except KeyError:
        print("❌ 'answer' column is missing.")
        sys.exit(1)

    allowed = {"a", "b", "c", "d", "A", "B", "C", "D"}
    invalid_answer_rows = [
        i + 2 for i, v in enumerate(answers)
        if str(v).strip() not in allowed
    ]

    if invalid_answer_rows:
        print("❌ File is not compatible for conversion.")
        print("   'answer' column must only contain a,b,c,d or A,B,C,D")
        print(f"   → Invalid values in rows: {invalid_answer_rows}")
        sys.exit(1)

    print("✅ Validation checks passed.\n")


# =========================================================
# LOADING EXCEL
# =========================================================

def read_questions_from_excel(file_path: str):
    """Reads CSV/Excel and runs validation."""
    if file_path.lower().endswith(".csv"):
        df = pd.read_csv(file_path, dtype=str)
    else:
        df = pd.read_excel(file_path, dtype=str)

    df.columns = [col.strip() for col in df.columns]

    required_cols = [
        "Question No", "question",
        "option-a", "option-b", "option-c", "option-d",
        "answer"
    ]
    missing = [c for c in required_cols if c not in df.columns]
    if missing:
        raise ValueError(f"Missing columns: {missing}")

    validate_excel(df)
    return df.fillna("")

#==============================================
# PPT HANDLER
#===============================================
def write_ppt(df, output_path):
    prs = Presentation()

    for idx, row in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout

        title = slide.shapes.title
        body = slide.placeholders[1]

        question = row["question"].strip()

        # -----------------------
        # Set Title
        # -----------------------
        title.text = question
        title_run = title.text_frame.paragraphs[0].runs[0]
        title_run.font.name = config.PPT_QUESTION_FONT
        title_run.font.size = Pt(config.PPT_QUESTION_FONT_SIZE)

        # -----------------------
        # Build options
        # -----------------------
        options = [
            row["option-a"].strip(),
            row["option-b"].strip(),
            row["option-c"].strip(),
            row["option-d"].strip(),
        ]

        bullets = ["a)", "b)", "c)", "d)"]

        tf = body.text_frame
        tf.text = ""   # clear placeholder

        for i, option in enumerate(options):
            p = tf.add_paragraph()

            # --- Remove built-in bullets ---
            p_pr = p._pPr
            buNone = OxmlElement('a:buNone')
            p_pr.insert(0, buNone)

            # --- Set custom text ---
            p.text = f"{bullets[i]} {option}"

            # Apply font to runs
            for run in p.runs:
                run.font.name = config.PPT_OPTION_FONT
                run.font.size = Pt(config.PPT_OPTION_FONT_SIZE)

        # Remove auto first empty paragraph
        if len(tf.paragraphs) > 0:
            tf.paragraphs[0].text = ""

    prs.save(output_path)
    print(f"✅ PPT file created: {os.path.abspath(output_path)}")

# =========================================================
# FORMATTING HELPERS
# =========================================================

def _set_run_font(run, name="Times New Roman", size_pt=11, bold=False):
    run.font.name = name
    try:
        run._element.rFonts.set(qn("w:eastAsia"), name)
    except:
        pass
    run.font.size = Pt(size_pt)
    run.font.bold = bold


def _set_paragraph_spacing(paragraph):
    paragraph.paragraph_format.line_spacing = config.LINE_SPACING


def set_section_margins(section):
    section.top_margin = Cm(config.TOP_MARGIN_CM)
    section.bottom_margin = Cm(config.BOTTOM_MARGIN_CM)
    section.left_margin = Cm(config.LEFT_MARGIN_CM)
    section.right_margin = Cm(config.RIGHT_MARGIN_CM)


def make_two_column_section(document):
    new_sec = document.add_section(WD_SECTION.CONTINUOUS)
    set_section_margins(new_sec)

    sectPr = new_sec._sectPr
    cols = OxmlElement("w:cols")
    cols.set(qn("w:num"), "2")
    cols.set(qn("w:sep"), "360")
    sectPr.append(cols)

    return new_sec


def make_single_column_section(document):
    new_sec = document.add_section(WD_SECTION.CONTINUOUS)
    set_section_margins(new_sec)

    sectPr = new_sec._sectPr
    cols = OxmlElement("w:cols")
    cols.set(qn("w:num"), "1")
    sectPr.append(cols)

    return new_sec


def create_doc_with_heading(document, heading_text):
    """Writes heading at top (full width)."""
    first_sec = document.sections[0]
    set_section_margins(first_sec)

    p = document.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(heading_text)
    _set_run_font(run, config.HEADING_FONT, config.HEADING_FONT_SIZE, bold=True)

    document.add_paragraph()  # spacing


def apply_layout_format(document):
    """Choose layout based on config.FORMAT."""
    layout_map = {
        1: make_single_column_section,
        2: make_two_column_section,
    }

    layout_fn = layout_map.get(config.FORMAT)

    if layout_fn:
        layout_fn(document)
    else:
        print(f"⚠️ Unknown FORMAT '{config.FORMAT}'. Defaulting to single column.")
        make_single_column_section(document)


# =========================================================
# WRITING QUESTIONS
# =========================================================

def write_questions_to_word(df, output_path):
    document = Document()
    create_doc_with_heading(document, config.TOPIC_NAME)
    apply_layout_format(document)

    errors = []

    for idx, row in df.iterrows():
        try:
            q_no = row["Question No"].strip()
            question = row["question"].strip()

            p_q = document.add_paragraph()
            run_q = p_q.add_run(f"{q_no}) {question}")
            _set_run_font(run_q, config.QUESTION_FONT, config.QUESTION_FONT_SIZE)
            _set_paragraph_spacing(p_q)

            p_opt = document.add_paragraph()
            p_opt.paragraph_format.left_indent = Pt(12)

            # Build the entire options line
            options_text = (
                f"a) {row['option-a']}    "
                f"b) {row['option-b']}    "
                f"c) {row['option-c']}    "
                f"d) {row['option-d']}"
            )

            run_opt = p_opt.add_run(options_text)
            _set_run_font(run_opt, config.OPTION_FONT, config.OPTION_FONT_SIZE)
            _set_paragraph_spacing(p_opt)

        except Exception as e:
            errors.append((idx + 1, str(e)))
            continue

    # Save file
    try:
        document.save(output_path)
    except PermissionError:
        alt = output_path.replace(".docx", "_new.docx")
        document.save(alt)
        print(f"⚠️ File open. Saved as {alt}")

    if errors:
        with open("errors.txt", "w") as f:
            for r, e in errors:
                f.write(f"Row {r}: {e}\n")

    print(f"✅ Questions file created: {os.path.abspath(output_path)}")


# =========================================================
# WRITING ANSWERS
# =========================================================

def write_answer_key_to_word(df, output_path):
    document = Document()
    first_section = document.sections[0]
    set_section_margins(first_section)

    p = document.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Answer Key")
    _set_run_font(run, config.HEADING_FONT, config.HEADING_FONT_SIZE, bold=True)
    _set_paragraph_spacing(p)

    document.add_paragraph()

    errors = []

    for idx, row in df.iterrows():
        try:
            q_no = row["Question No"].strip()
            ans_val = row["answer"].strip()
            ans_lower = ans_val.lower()

            option_map = {
                "a": row["option-a"].strip(),
                "b": row["option-b"].strip(),
                "c": row["option-c"].strip(),
                "d": row["option-d"].strip(),
            }

            actual_answer_text = option_map.get(ans_lower, "Invalid Option")

            p_ans = document.add_paragraph()
            run_ans = p_ans.add_run(f"Q{q_no}) {ans_val} -> {actual_answer_text}")
            _set_run_font(run_ans, config.OPTION_FONT, config.OPTION_FONT_SIZE)
            _set_paragraph_spacing(p_ans)

        except Exception as e:
            errors.append((idx + 1, str(e)))
            continue

    try:
        document.save(output_path)
    except PermissionError:
        alt = output_path.replace(".docx", "_new.docx")
        document.save(alt)
        print(f"⚠️ File open. Saved as {alt}")

    print(f"✅ Answer key file created: {os.path.abspath(output_path)}")


# =========================================================
# MAIN
# =========================================================

def main():
    parser = argparse.ArgumentParser(description="Generate Word and optionally PPT from questions file.")
    parser.add_argument("--ppt", nargs="?", const="y", help="Generate PPT file if provided")
    args = parser.parse_args()

    os.makedirs(config.INPUT_FOLDER, exist_ok=True)
    os.makedirs(config.OUTPUT_FOLDER, exist_ok=True)

    df = read_questions_from_excel(config.INPUT_FILE)

    write_questions_to_word(df, config.OUTPUT_QUESTIONS_FILE)
    write_answer_key_to_word(df, config.OUTPUT_ANSWERS_FILE)

    # If --ppt is provided, generate PPT
    if args.ppt:
        write_ppt(df, config.OUTPUT_PPT_FILE)

# 
# def main():
#     os.makedirs(config.INPUT_FOLDER, exist_ok=True)
#     os.makedirs(config.OUTPUT_FOLDER, exist_ok=True)
#     df = read_questions_from_excel(config.INPUT_FILE)
#     write_questions_to_word(df, config.OUTPUT_QUESTIONS_FILE)
#     write_answer_key_to_word(df, config.OUTPUT_ANSWERS_FILE)


if __name__ == "__main__":
    main()